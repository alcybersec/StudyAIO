"""PPTX extractor using python-pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import structlog

from app.core.exceptions import ExtractionError
from app.extractors.base import (
    BaseExtractor,
    ExtractionResult,
    ImageInfo,
    PageContent,
)

logger = structlog.get_logger()


class PptxExtractor(BaseExtractor):
    """Extracts text and images from PPTX files using python-pptx."""

    def extract(self, file_path: Path, output_dir: Path) -> ExtractionResult:
        """Extract text and images from a PPTX file.

        Each slide maps to a page in the extraction result. Speaker notes
        are appended to the slide text.

        Args:
            file_path: Path to the PPTX file.
            output_dir: Directory to save extracted images.

        Returns:
            ExtractionResult with per-slide content.

        Raises:
            ExtractionError: If the PPTX cannot be opened or parsed.
        """
        images_dir = output_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            raise ExtractionError(f"Failed to open PPTX {file_path.name}: {e}") from e

        pages: list[PageContent] = []
        total_images = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_number = slide_idx + 1
            text_parts: list[str] = []
            slide_images: list[ImageInfo] = []

            for shape in slide.shapes:
                # Extract text from text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        ]
                        if row_texts:
                            text_parts.append(" | ".join(row_texts))

                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        content_type = image.content_type
                        ext = content_type.split("/")[-1].replace("jpeg", "jpg")
                        total_images += 1
                        filename = f"slide{slide_number}_img{total_images}.{ext}"
                        image_path = images_dir / filename
                        image_path.write_bytes(image.blob)

                        slide_images.append(
                            ImageInfo(
                                filename=filename,
                                position=f"slide_{slide_number}",
                            )
                        )
                    except Exception:
                        logger.warning(
                            "pptx_image_extraction_failed",
                            slide=slide_number,
                        )

            # Append speaker notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_parts.append(f"\n[Speaker Notes]\n{notes_text}")

            pages.append(
                PageContent(
                    page_number=slide_number,
                    text="\n".join(text_parts),
                    images=slide_images,
                )
            )

        logger.info(
            "pptx_extraction_complete",
            file=file_path.name,
            slides=len(pages),
            images=total_images,
        )

        return ExtractionResult(
            pages=pages,
            metadata={"extractor_version": "1.0", "source_type": "pptx"},
            image_count=total_images,
            page_count=len(pages),
        )
