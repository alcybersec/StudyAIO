"""Pipeline stage 1: Classify — identify course, week, and title."""

import asyncio
from datetime import datetime
from pathlib import Path

import structlog
from sqlalchemy import select

from app.agents.factory import get_agent
from app.config import settings
from app.core.database import async_session_factory
from app.core.exceptions import AgentError, ClassificationError
from app.core.utils import generate_id
from app.models.artifact import LectureArtifact
from app.models.course import Course
from app.models.pipeline_run import PipelineRun
from app.services import review_service
from app.worker import celery_app

logger = structlog.get_logger()


def _run_async(coro):
    """Run an async coroutine from a sync Celery task."""
    loop = asyncio.new_event_loop()
    try:
        return loop.run_until_complete(coro)
    finally:
        loop.close()


def _extract_text_preview(file_path: str, file_type: str) -> str:
    """Extract text from the first 2 pages for classification.

    Uses a lightweight extraction — just text, no images.

    Args:
        file_path: Path to the file.
        file_type: One of "pdf", "docx", "pptx".

    Returns:
        Text preview string.
    """
    path = Path(file_path)
    if not path.exists():
        return ""

    try:
        if file_type == "pdf":
            import fitz
            doc = fitz.open(str(path))
            texts = []
            for i in range(min(2, len(doc))):
                texts.append(doc[i].get_text("text"))
            doc.close()
            return "\n\n".join(texts)

        elif file_type == "docx":
            import docx
            document = docx.Document(str(path))
            texts = []
            for para in document.paragraphs[:50]:
                if para.text.strip():
                    texts.append(para.text)
            return "\n".join(texts[:50])

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text)
            return "\n".join(texts)

    except Exception as e:
        logger.warning("text_preview_extraction_failed", error=str(e))
        return ""

    return ""


async def _get_known_courses(session) -> list[str]:
    """Get all existing course codes."""
    result = await session.execute(select(Course.code))
    return [row[0] for row in result.all()]


async def _get_or_create_course(session, course_code: str) -> Course:
    """Get existing course or create a new one.

    Args:
        session: Database session.
        course_code: Course code to find or create.

    Returns:
        Course instance.
    """
    result = await session.execute(
        select(Course).where(Course.code == course_code)
    )
    course = result.scalar_one_or_none()
    if course:
        return course

    course = Course(
        id=generate_id(),
        code=course_code,
    )
    session.add(course)
    await session.flush()
    logger.info("course_created", course_id=course.id, code=course.code)
    return course


async def _classify(artifact_id: str) -> dict:
    """Async classify implementation."""
    async with async_session_factory() as session:
        # Load artifact
        result = await session.execute(
            select(LectureArtifact).where(LectureArtifact.id == artifact_id)
        )
        artifact = result.scalar_one_or_none()
        if not artifact:
            raise ClassificationError(f"Artifact {artifact_id} not found")

        # Update status
        artifact.status = "classifying"
        await session.commit()

        # Create pipeline run
        run = PipelineRun(
            id=generate_id(),
            artifact_id=artifact_id,
            stage="classify",
            status="running",
            started_at=datetime.utcnow(),
        )
        session.add(run)
        await session.flush()

        try:
            # Extract text preview for classification
            text_preview = _extract_text_preview(artifact.file_path, artifact.file_type)
            if not text_preview:
                raise ClassificationError(
                    f"Could not extract text preview from {artifact.original_filename}"
                )

            # Get known courses
            known_courses = await _get_known_courses(session)

            # Call AI agent
            agent = get_agent()
            classification = await agent.classify_lecture(
                text_preview=text_preview,
                filename=artifact.original_filename,
                known_courses=known_courses,
            )

            logger.info(
                "classification_result",
                artifact_id=artifact_id,
                course_code=classification.course_code,
                week=classification.week,
                confidence=classification.confidence,
            )

            threshold = settings.classification_confidence_threshold

            if classification.confidence >= threshold:
                # High confidence — apply classification
                course = await _get_or_create_course(
                    session, classification.course_code
                )
                artifact.course_id = course.id
                artifact.week = classification.week
                artifact.title = classification.title
                artifact.status = "classified"

                run.status = "completed"
                run.completed_at = datetime.utcnow()
                if run.started_at:
                    delta = run.completed_at - run.started_at
                    run.duration_ms = int(delta.total_seconds() * 1000)

                await session.commit()

                return {
                    "artifact_id": artifact_id,
                    "status": "classified",
                    "course_code": classification.course_code,
                    "week": classification.week,
                    "title": classification.title,
                    "confidence": classification.confidence,
                }

            else:
                # Low confidence — create review item
                review_type = "classification_course"
                if classification.course_code != "UNKNOWN" and classification.week == 0:
                    review_type = "classification_week"

                # Build payload per PRD spec
                payload = {
                    "context": text_preview[:2000],
                    "filename": artifact.original_filename,
                    "suggestions": [
                        {
                            "value": {
                                "course": classification.course_code,
                                "week": classification.week,
                            },
                            "confidence": classification.confidence,
                        }
                    ],
                    "reason": classification.reasoning
                    or "Classification confidence below threshold",
                }

                suggested_values = {
                    "course_code": classification.course_code,
                    "week": classification.week,
                    "title": classification.title,
                }

                await review_service.create_review_item(
                    session=session,
                    review_type=review_type,
                    entity_type="lecture_artifact",
                    entity_id=artifact_id,
                    payload=payload,
                    suggested_values=suggested_values,
                )

                artifact.status = "waiting_review"
                run.status = "waiting_review"
                run.completed_at = datetime.utcnow()
                if run.started_at:
                    delta = run.completed_at - run.started_at
                    run.duration_ms = int(delta.total_seconds() * 1000)

                await session.commit()

                return {
                    "artifact_id": artifact_id,
                    "status": "waiting_review",
                    "confidence": classification.confidence,
                    "review_type": review_type,
                }

        except (ClassificationError, AgentError):
            run.status = "failed"
            run.error_message = str(artifact_id)
            run.completed_at = datetime.utcnow()
            artifact.status = "failed"
            await session.commit()
            raise

        except Exception as e:
            run.status = "failed"
            run.error_message = str(e)
            run.completed_at = datetime.utcnow()
            artifact.status = "failed"
            await session.commit()
            raise ClassificationError(f"Classification failed: {e}") from e


@celery_app.task(
    name="app.pipeline.classify.classify_artifact",
    bind=True,
    max_retries=2,
    default_retry_delay=15,
)
def classify_artifact(self, input_value: str | dict) -> dict:
    """Celery task: classify an ingested artifact.

    Accepts either a plain artifact_id string or a dict from the
    previous pipeline stage (for chain compatibility).

    Args:
        input_value: Artifact UUID string or dict with artifact_id.

    Returns:
        Dict with artifact_id, status, course_code, week, title, confidence.
    """
    # Resolve input (chain compatibility)
    if isinstance(input_value, dict):
        status = input_value.get("status", "")
        if status in ("duplicate", "waiting_review", "failed"):
            logger.info(
                "classify_task_skipped",
                status=status,
                artifact_id=input_value.get("artifact_id"),
            )
            return input_value
        artifact_id = input_value.get("artifact_id", "")
    else:
        artifact_id = input_value

    if not artifact_id:
        raise ClassificationError("No artifact_id provided")

    logger.info("classify_task_started", artifact_id=artifact_id)
    try:
        return _run_async(_classify(artifact_id))
    except (ClassificationError, AgentError):
        raise  # Don't retry on classification/agent errors
    except Exception as exc:
        logger.error("classify_task_error", error=str(exc), artifact_id=artifact_id)
        raise self.retry(exc=exc)
