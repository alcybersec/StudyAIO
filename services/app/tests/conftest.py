"""Shared test fixtures for StudyAIO."""

import tempfile
from unittest.mock import AsyncMock, MagicMock, patch

import httpx
import pytest

from app.agents.base import (
    ClassificationResult,
    SummaryResult,
)
from app.core.storage import reset_storage

# ── Sample data dicts ──────────────────────────────────────────────


@pytest.fixture
def sample_course_data():
    """Minimal course data dict."""
    return {
        "id": "course-001",
        "code": "CSIT302",
        "name": "Cybersecurity",
    }


@pytest.fixture
def sample_artifact_data():
    """Minimal artifact data dict."""
    return {
        "id": "artifact-001",
        "course_id": "course-001",
        "week": 5,
        "title": "Network Security",
        "original_filename": "CSIT302_Week5.pdf",
        "file_path": "/app/data/uploads/artifact-001_CSIT302_Week5.pdf",
        "file_type": "pdf",
        "sha256": "a" * 64,
        "file_size_bytes": 1024,
        "status": "ingested",
    }


@pytest.fixture
def sample_manifest():
    """Sample extraction manifest dict."""
    return {
        "pages": [
            {
                "page_number": 1,
                "text": "Introduction to Network Security\nCSIT302 Week 5",
                "images": [{"filename": "page1_img1.png", "caption": "", "position": "page_1"}],
            },
            {
                "page_number": 2,
                "text": "Firewall types and configurations\nPacket filtering vs stateful inspection",
                "images": [],
            },
        ],
        "metadata": {"extractor_version": "1.0", "source_type": "pdf"},
    }


# ── Mock agent ──────────────────────────────────────────────────────


@pytest.fixture
def mock_agent():
    """AsyncMock of AgentAdapter with sensible return values."""
    agent = AsyncMock()
    agent.classify_lecture.return_value = ClassificationResult(
        course_code="CSIT302",
        week=5,
        title="Network Security",
        confidence=0.92,
        reasoning="Course code in header, week in filename",
    )
    agent.generate_summary.return_value = SummaryResult(
        content_md="# CSIT302 — Week 5: Network Security\n\n## Key Concepts\n- Firewalls",
        embedded_images=[],
    )
    return agent


# ── Mock DB session ─────────────────────────────────────────────────


@pytest.fixture
def mock_session():
    """AsyncMock of AsyncSession with basic query support.

    The execute() return value is a MagicMock (not AsyncMock) so that
    synchronous SQLAlchemy Result methods like .scalars().all() and
    .scalar_one_or_none() work correctly without returning coroutines.
    """
    session = AsyncMock()
    session.add = MagicMock()
    session.flush = AsyncMock()
    session.commit = AsyncMock()
    session.refresh = AsyncMock()
    session.execute = AsyncMock()
    # Return a MagicMock for the execute result so sync chains work
    execute_result = MagicMock()
    execute_result.scalars.return_value.all.return_value = []
    execute_result.scalars.return_value.first.return_value = None
    execute_result.scalar_one_or_none.return_value = None
    execute_result.all.return_value = []
    session.execute.return_value = execute_result
    return session


# ── Async HTTP test client ────────────────────────────────────────


@pytest.fixture
def default_test_user(make_user):
    """Default user for API tests (simulates get_current_user_or_default)."""
    return make_user(
        id="00000000-0000-0000-0000-000000000001",
        email="admin@studyaio.local",
        username="admin",
        role="admin",
        tier="pro",
    )


@pytest.fixture
async def async_client(mock_session, default_test_user):
    """Async HTTP client for testing FastAPI endpoints."""
    from app.api.deps import get_current_user_or_default
    from app.core.database import get_session
    from app.main import app

    async def override_session():
        yield mock_session

    async def override_user():
        return default_test_user

    app.dependency_overrides[get_session] = override_session
    app.dependency_overrides[get_current_user_or_default] = override_user

    # Reset rate limiter state between tests to prevent cross-test 429s
    from app.core.rate_limit import limiter

    limiter.reset()

    # Use a writable temp dir for data_dir so upload tests work as non-root
    with tempfile.TemporaryDirectory() as tmpdir, patch("app.config.settings.data_dir", tmpdir):
        # Reset storage singleton so it picks up the patched data_dir
        reset_storage()
        async with httpx.AsyncClient(
            transport=httpx.ASGITransport(app=app),
            base_url="http://test",
        ) as client:
            yield client
        # Reset again on teardown so other tests get a fresh singleton
        reset_storage()
    app.dependency_overrides.clear()


# ── Auth helpers ──────────────────────────────────────────────────────


@pytest.fixture
def make_user():
    """Factory fixture to create mock User objects for testing auth."""
    from app.core.auth import hash_password
    from app.models.user import User

    def _make(
        id: str = "user-001",
        email: str = "test@example.com",
        username: str = "testuser",
        password: str = "TestPass1!",
        role: str = "user",
        tier: str = "free",
        is_active: bool = True,
        mfa_enabled: bool = False,
        mfa_secret: str | None = None,
        **kwargs,
    ) -> User:
        user = MagicMock(spec=User)
        user.id = id
        user.email = email
        user.username = username
        user.hashed_password = hash_password(password)
        user.role = role
        user.tier = tier
        user.is_active = is_active
        user.email_verified = False
        user.mfa_enabled = mfa_enabled
        user.mfa_secret = mfa_secret
        user.avatar_url = None
        user.backup_codes = None
        user.last_login_at = None
        # None = no token cutoff; pass tokens_valid_from=... to simulate a
        # password reset/change having revoked earlier tokens.
        user.tokens_valid_from = None
        user.created_at = MagicMock()
        user.updated_at = MagicMock()
        for k, v in kwargs.items():
            setattr(user, k, v)
        return user

    return _make


@pytest.fixture
def auth_cookies(make_user):
    """Helper to generate auth cookies for a user."""
    from app.core.auth import create_access_token, create_refresh_token

    def _cookies(user=None, **user_kwargs):
        if user is None:
            user = make_user(**user_kwargs)
        access = create_access_token(user.id, user.role, user.tier)
        refresh = create_refresh_token(user.id)
        return {"access_token": access, "refresh_token": refresh}, user

    return _cookies


# ── Programmatic fixture file creators ──────────────────────────────


@pytest.fixture
def simple_pdf(tmp_path):
    """Create a minimal PDF file using PyMuPDF."""
    import fitz

    pdf_path = tmp_path / "test_lecture.pdf"
    doc = fitz.open()

    # Page 1
    page1 = doc.new_page(width=595, height=842)
    page1.insert_text((72, 72), "CSIT302 Week 5\nNetwork Security Fundamentals", fontsize=14)

    # Page 2
    page2 = doc.new_page(width=595, height=842)
    page2.insert_text(
        (72, 72), "Firewalls are network security systems\nthat monitor traffic.", fontsize=12
    )

    doc.save(str(pdf_path))
    doc.close()
    return pdf_path


@pytest.fixture
def simple_docx(tmp_path):
    """Create a minimal DOCX file using python-docx."""
    import docx

    docx_path = tmp_path / "test_lecture.docx"
    document = docx.Document()

    document.add_heading("CSIT302 Week 5: Network Security", level=1)
    document.add_paragraph("This lecture covers network security fundamentals.")
    document.add_heading("Firewalls", level=2)
    document.add_paragraph(
        "A firewall monitors and controls incoming and outgoing network traffic."
    )

    document.save(str(docx_path))
    return docx_path


@pytest.fixture
def simple_pptx(tmp_path):
    """Create a minimal PPTX file using python-pptx."""
    from pptx import Presentation

    pptx_path = tmp_path / "test_lecture.pptx"
    prs = Presentation()

    # Slide 1 — title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CSIT302 Week 5"
    slide1.placeholders[1].text = "Network Security Fundamentals"

    # Slide 2 — content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Firewalls"
    slide2.placeholders[1].text = "Packet filtering\nStateful inspection\nApplication layer"

    # Slide 3 — with speaker notes
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "IDS/IPS"
    slide3.placeholders[1].text = "Intrusion Detection and Prevention Systems"
    notes_slide = slide3.notes_slide
    notes_slide.notes_text_frame.text = "Important: Explain difference between IDS and IPS"

    prs.save(str(pptx_path))
    return pptx_path
