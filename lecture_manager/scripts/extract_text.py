#!/usr/bin/env python3
"""
Extract text and images from lecture files (.docx, .pdf, .pptx).

Usage:
    python extract_text.py <filepath>
    python extract_text.py <filepath> --first-n-chars 2000
    python extract_text.py <filepath> --extract-images <output_dir>

Text is printed to stdout.
Images are saved to <output_dir>/ as <basename>_img001.png, _img002.jpg, etc.
A manifest of extracted images is printed to stderr for reference.

Dependencies (install if missing):
    pip install pypdf pdfplumber python-docx markitdown[pptx] python-pptx Pillow --break-system-packages
"""

import sys
import os
import argparse
import zipfile
import hashlib
from pathlib import Path


# ---------------------------------------------------------------------------
# Image extraction helpers
# ---------------------------------------------------------------------------

SIGNATURE_MAP = {
    b'\x89PNG': '.png',
    b'\xff\xd8\xff': '.jpg',
    b'GIF87a': '.gif',
    b'GIF89a': '.gif',
    b'RIFF': '.webp',
    b'BM': '.bmp',
    b'<svg': '.svg',
    b'II*\x00': '.tiff',
    b'MM\x00*': '.tiff',
}


def _guess_ext(data: bytes) -> str:
    """Guess image extension from raw bytes."""
    for sig, ext in SIGNATURE_MAP.items():
        if data[:len(sig)] == sig:
            return ext
    return '.png'


def _image_hash(data: bytes) -> str:
    return hashlib.md5(data).hexdigest()


def _save_image(data: bytes, output_dir: str, basename: str, counter: int,
                seen_hashes: set, context: str = "") -> tuple:
    """
    Save image bytes to output_dir if not a duplicate and not tiny.
    Normalizes all images to PNG or JPG for API compatibility.
    Resizes images larger than MAX_DIM to prevent processing errors.
    Returns (new_counter, saved_filepath_or_None).
    """
    MAX_DIM = 2048  # max width or height in pixels
    MAX_BYTES = 5 * 1024 * 1024  # 5 MB cap for output file

    # Skip tiny images (icons/bullets, < 2 KB)
    if len(data) < 2048:
        return counter, None

    h = _image_hash(data)
    if h in seen_hashes:
        return counter, None
    seen_hashes.add(h)

    # Normalize with Pillow: convert to standard PNG/JPG and resize if needed
    try:
        from PIL import Image
        import io

        img = Image.open(io.BytesIO(data))

        # Convert palette/CMYK/RGBA modes to RGB for broad compatibility
        if img.mode in ("P", "CMYK"):
            img = img.convert("RGB")
        elif img.mode == "RGBA":
            # Keep RGBA for PNGs (transparency), convert for JPEGs
            pass
        elif img.mode not in ("RGB", "L", "RGBA"):
            img = img.convert("RGB")

        # Resize if too large
        w, h_px = img.size
        if w > MAX_DIM or h_px > MAX_DIM:
            ratio = min(MAX_DIM / w, MAX_DIM / h_px)
            new_size = (int(w * ratio), int(h_px * ratio))
            img = img.resize(new_size, Image.LANCZOS)

        # Save as PNG (lossless, handles transparency) or JPG (smaller for photos)
        buf = io.BytesIO()
        if img.mode == "RGBA":
            img.save(buf, format="PNG", optimize=True)
            ext = ".png"
        else:
            # Try JPEG first (smaller), fall back to PNG
            img.save(buf, format="JPEG", quality=85, optimize=True)
            ext = ".jpg"

        out_data = buf.getvalue()

        # If still too big, reduce quality further
        if len(out_data) > MAX_BYTES and ext == ".jpg":
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=60, optimize=True)
            out_data = buf.getvalue()

        filename = f"{basename}_img{counter:03d}{ext}"
        filepath = os.path.join(output_dir, filename)
        with open(filepath, 'wb') as f:
            f.write(out_data)

        ctx_note = f" ({context})" if context else ""
        dims = f" [{img.size[0]}x{img.size[1]}]"
        print(f"[IMAGE] Saved: {filename}{dims}{ctx_note}", file=sys.stderr)
        return counter + 1, filepath

    except Exception as e:
        # If Pillow can't handle it, skip rather than save a broken file
        print(f"[WARNING] Could not normalize image (counter={counter}): {e}", file=sys.stderr)
        return counter, None


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def extract_images_from_pdf(filepath: str, output_dir: str, basename: str,
                            counter: int, seen_hashes: set) -> int:
    """Extract embedded raster images from PDF, then render pages as fallback."""
    # 1. Embedded images via pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        for page_num, page in enumerate(reader.pages):
            resources = page.get("/Resources")
            if not resources:
                continue
            x_objects = resources.get("/XObject")
            if not x_objects:
                continue
            x_objects = x_objects.get_object()
            for obj_name in x_objects:
                obj = x_objects[obj_name].get_object()
                if obj.get("/Subtype") == "/Image":
                    try:
                        data = obj.get_data()
                        if data:
                            counter, _ = _save_image(
                                data, output_dir, basename, counter, seen_hashes,
                                context=f"PDF page {page_num + 1}"
                            )
                    except Exception:
                        pass
    except Exception as e:
        print(f"[WARNING] PDF embedded image extraction: {e}", file=sys.stderr)

    # 2. Render full pages (catches vector diagrams, charts, etc.)
    try:
        import subprocess
        import tempfile
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["pdftoppm", "-jpeg", "-r", "200", filepath, os.path.join(tmpdir, "page")],
                capture_output=True, timeout=120
            )
            if result.returncode == 0:
                for page_img in sorted(Path(tmpdir).glob("page-*.jpg")):
                    data = page_img.read_bytes()
                    page_label = page_img.stem.replace("page-", "page")
                    counter, _ = _save_image(
                        data, output_dir, basename, counter, seen_hashes,
                        context=f"PDF {page_label} render"
                    )
    except FileNotFoundError:
        print("[INFO] pdftoppm not found — skipping full page renders", file=sys.stderr)
    except Exception as e:
        print(f"[WARNING] PDF page rendering: {e}", file=sys.stderr)

    return counter


def extract_from_pdf(filepath: str) -> str:
    text = ""
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        if text.strip():
            return text.strip()
    except ImportError:
        pass
    except Exception as e:
        print(f"[WARNING] pdfplumber failed: {e}, trying pypdf...", file=sys.stderr)

    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n\n"
        return text.strip()
    except Exception as e:
        print(f"[ERROR] PDF text extraction failed: {e}", file=sys.stderr)
        return ""


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_images_from_docx(filepath: str, output_dir: str, basename: str,
                             counter: int, seen_hashes: set) -> int:
    """Extract images from word/media/ inside the DOCX zip."""
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            media_files = [f for f in z.namelist() if f.startswith('word/media/')]
            for media_path in sorted(media_files):
                data = z.read(media_path)
                original_name = os.path.basename(media_path)
                counter, _ = _save_image(
                    data, output_dir, basename, counter, seen_hashes,
                    context=f"DOCX embedded: {original_name}"
                )
    except Exception as e:
        print(f"[WARNING] DOCX image extraction: {e}", file=sys.stderr)
    return counter


def extract_from_docx(filepath: str) -> str:
    text = ""
    try:
        import docx
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                text += row_text + "\n"
            text += "\n"
        if text.strip():
            return text.strip()
    except ImportError:
        pass
    except Exception as e:
        print(f"[WARNING] python-docx failed: {e}, trying pandoc...", file=sys.stderr)

    try:
        import subprocess
        result = subprocess.run(
            ["pandoc", filepath, "-t", "plain", "--wrap=none"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            return result.stdout.strip()
    except Exception as e:
        print(f"[ERROR] DOCX text extraction failed: {e}", file=sys.stderr)
    return ""


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_images_from_pptx(filepath: str, output_dir: str, basename: str,
                             counter: int, seen_hashes: set) -> int:
    """Extract images from ppt/media/ inside the PPTX zip, plus render slides."""
    # 1. Embedded media files
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
            for media_path in sorted(media_files):
                data = z.read(media_path)
                original_name = os.path.basename(media_path)
                counter, _ = _save_image(
                    data, output_dir, basename, counter, seen_hashes,
                    context=f"PPTX embedded: {original_name}"
                )
    except Exception as e:
        print(f"[WARNING] PPTX embedded image extraction: {e}", file=sys.stderr)

    # 2. Render slides as images (captures SmartArt, charts, drawn diagrams)
    try:
        import subprocess
        import tempfile
        with tempfile.TemporaryDirectory() as tmpdir:
            # Convert PPTX → PDF via LibreOffice
            soffice_script = os.path.join(
                os.path.dirname(os.path.abspath(__file__)), "office", "soffice.py"
            )
            if os.path.exists(soffice_script):
                subprocess.run(
                    ["python", soffice_script, "--headless",
                     "--convert-to", "pdf", "--outdir", tmpdir, filepath],
                    capture_output=True, timeout=120
                )
            else:
                # Try system LibreOffice directly
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pdf",
                     "--outdir", tmpdir, filepath],
                    capture_output=True, timeout=120
                )

            pdfs = list(Path(tmpdir).glob("*.pdf"))
            if pdfs:
                subprocess.run(
                    ["pdftoppm", "-jpeg", "-r", "200", str(pdfs[0]),
                     os.path.join(tmpdir, "slide")],
                    capture_output=True, timeout=120
                )
                for slide_img in sorted(Path(tmpdir).glob("slide-*.jpg")):
                    data = slide_img.read_bytes()
                    slide_label = slide_img.stem.replace("slide-", "slide")
                    counter, _ = _save_image(
                        data, output_dir, basename, counter, seen_hashes,
                        context=f"PPTX {slide_label} render"
                    )
    except FileNotFoundError:
        print("[INFO] LibreOffice/pdftoppm not available — skipping slide renders", file=sys.stderr)
    except Exception as e:
        print(f"[WARNING] PPTX slide rendering: {e}", file=sys.stderr)

    return counter


def extract_from_pptx(filepath: str) -> str:
    text = ""
    try:
        import subprocess
        result = subprocess.run(
            ["python", "-m", "markitdown", filepath],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except Exception as e:
        print(f"[WARNING] markitdown failed: {e}, trying python-pptx...", file=sys.stderr)

    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            text += f"--- Slide {i + 1} ---\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text += para.text + "\n"
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        text += row_text + "\n"
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text += f"[Speaker Notes: {notes}]\n"
            text += "\n"
        return text.strip()
    except ImportError:
        pass
    except Exception as e:
        print(f"[ERROR] PPTX text extraction failed: {e}", file=sys.stderr)
    return ""


# ---------------------------------------------------------------------------
# Unified interface
# ---------------------------------------------------------------------------

def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_from_pdf(filepath)
    elif ext == ".docx":
        return extract_from_docx(filepath)
    elif ext in (".pptx", ".ppt"):
        return extract_from_pptx(filepath)
    else:
        print(f"[ERROR] Unsupported file type: {ext}", file=sys.stderr)
        return ""


def extract_images(filepath: str, output_dir: str) -> list:
    """
    Extract all images from a file into output_dir.
    Returns list of saved image filepaths.
    """
    os.makedirs(output_dir, exist_ok=True)
    basename = Path(filepath).stem
    ext = os.path.splitext(filepath)[1].lower()
    seen_hashes = set()
    counter = 1

    if ext == ".pdf":
        counter = extract_images_from_pdf(filepath, output_dir, basename, counter, seen_hashes)
    elif ext == ".docx":
        counter = extract_images_from_docx(filepath, output_dir, basename, counter, seen_hashes)
    elif ext in (".pptx", ".ppt"):
        counter = extract_images_from_pptx(filepath, output_dir, basename, counter, seen_hashes)
    else:
        print(f"[ERROR] Unsupported file type: {ext}", file=sys.stderr)

    saved = sorted(Path(output_dir).glob(f"{basename}_img*"))
    return [str(p) for p in saved]


def main():
    parser = argparse.ArgumentParser(
        description="Extract text and images from lecture files (.docx, .pdf, .pptx)"
    )
    parser.add_argument("filepath", help="Path to the lecture file")
    parser.add_argument("--first-n-chars", type=int, default=0,
                        help="Only output the first N characters of text (0 = all)")
    parser.add_argument("--extract-images", metavar="OUTPUT_DIR", default=None,
                        help="Extract images to this directory")
    args = parser.parse_args()

    if not os.path.exists(args.filepath):
        print(f"[ERROR] File not found: {args.filepath}", file=sys.stderr)
        sys.exit(1)

    # Extract text
    text = extract_text(args.filepath)
    if not text:
        print("[WARNING] No text could be extracted.", file=sys.stderr)
    else:
        if args.first_n_chars > 0:
            text = text[:args.first_n_chars]
        print(text)

    # Extract images if requested
    if args.extract_images:
        images = extract_images(args.filepath, args.extract_images)
        if images:
            print(f"\n[IMAGE MANIFEST] {len(images)} image(s) extracted:", file=sys.stderr)
            for img in images:
                print(f"  - {img}", file=sys.stderr)
        else:
            print("[INFO] No images found in this file.", file=sys.stderr)


if __name__ == "__main__":
    main()
